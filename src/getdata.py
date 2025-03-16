import os
import docx
import pandas as pd
from docx.document import Document as DocxDocument
from docx.oxml.ns import qn
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from typing import List, Dict
import fitz 
from pptx import Presentation
import camelot
from agents.workflow import ProjectScopingAgent
from contextlib import asynccontextmanager
from io import BytesIO
import pdfplumber
import os
import tempfile
from utils.logger import logger


class ExtractText:
    def __init__(self, document_path:str = None, url=None,user_id:str=None, document_id:str=None):
        self.document_path = document_path
        self.url = url
        self.user_id = user_id
        self.document_id = document_id
        os.makedirs("uploads_images", exist_ok=True)
        logger.info(f"document_path: {self.document_path},user_id: {self.user_id}, document_id: {self.document_id}")

    @staticmethod
    def iter_block_items(parent):
        """
        Iterate through paragraphs, tables, and other block elements in document order.
        """

        if isinstance(parent, DocxDocument):
            parent_elm = parent.element.body
        else:
            parent_elm = parent

        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield docx.text.paragraph.Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield docx.table.Table(child, parent)


    async def extract_docx(self):
        doc = docx.Document(self.document_path)
        content = []
        image_count = 0
        os.makedirs("uploads_images", exist_ok=True)

        # Iterate through all elements in the document
        for block in self.__class__.iter_block_items(doc):
            # Process paragraphs
            if isinstance(block, docx.text.paragraph.Paragraph):
                text = block.text.strip()
                if text:
                    content.append({"type": "text", "data": text})

                # Extract images from the paragraph's XML
                for elem in block._element.iter():
                    if elem.tag.endswith('drawing'):
                        # Check for inline images
                        inline = elem.find('.//' + qn('wp:inline')) 
                        if inline is not None:
                            blip = inline.find('.//' + qn('a:blip'))
                            if blip is not None:
                                image_id = blip.get(qn('r:embed'))
                                if image_id:
                                    image_part = doc.part.related_parts[image_id]
                                    image_bytes = image_part.blob
                                    image_count += 1
                                    image_path = os.path.join("uploads_images", f"{self.document_id}_{self.user_id}_image_{image_count}.png")
                                    with open(image_path, "wb") as f:
                                        f.write(image_bytes)
                                    content.append({"type": "image", "data": image_path, "content":ProjectScopingAgent.summarize_image(image_path)})

                        # Check for floating images (anchored)
                        anchor = elem.find('.//' + qn('wp:anchor'))
                        if anchor is not None:
                            blip = anchor.find('.//' + qn('a:blip'))
                            if blip is not None:
                                image_id = blip.get(qn('r:embed'))
                                if image_id:
                                    image_part = doc.part.related_parts[image_id]
                                    image_bytes = image_part.blob
                                    image_count += 1
                                    image_path = os.path.join("uploads_images", f"{self.document_id}_{self.user_id}_image_{image_count}.png")
                                    with open(image_path, "wb") as f:
                                        f.write(image_bytes)
                                    content.append({"type": "image", "data": image_path, "content":ProjectScopingAgent.summarize_image(image_path)})

            # Process tables
            elif isinstance(block, docx.table.Table):
                table_data = []
                for row in block.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                content.append({"type": "table", "data": table_data})

        return content




    # async def process_pdf_with_structure(self):
    #     content = []

    #     # Open the PDF file directly without loading into memory
    #     with fitz.open(self.document_path) as doc:
    #         for page_num, page in enumerate(doc):
    #             # Optimized text extraction
    #             text = page.get_text("text").strip()
    #             if text:
    #                 content.append({"type": "text", "data": text})

    #             # Optimized image extraction
    #             for img_index, img in enumerate(page.get_images(full=True)):
    #                 xref = img[0]
    #                 pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # High-quality image
    #                 image_path = f"uploads_images/{self.document_id}_{self.user_id}_pdf_image_{page_num+1}_{img_index+1}.png"
    #                 pix.save(image_path)
    #                 content.append({"type": "image", "data": image_path, "content": "Image extracted"})

    #     # Optimized Table Extraction using pdfplumber
    #     with pdfplumber.open(self.document_path) as pdf:
    #         for page_num, page in enumerate(pdf.pages):
    #             tables = page.extract_tables()
    #             for table in tables:
    #                 content.append({"type": "table", "data": table})

    #     return content


    async def process_pdf_with_structure(self):
        content = []

        # Read the entire PDF into memory
        with open(self.document_path, "rb") as f:
            pdf_bytes = f.read()
            logger.info("completed reading the document uploaded")

        # Process text and images with PyMuPDF from memory
        async with open_pdf(stream=pdf_bytes, filetype="pdf") as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    content.append({"type": "text", "data": text})

                # Extract images
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    image_path = os.path.join("uploads_images", f"{self.document_id}_{self.user_id}_pdf_image_{page_num+1}_{img_index+1}.{image_ext}")
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    content.append({"type": "image", "data": image_path, "content": ProjectScopingAgent.summarize_image(image_path)})

        # Process tables with Camelot using a temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
        try:
            temp_file.write(pdf_bytes)
            temp_file.close()  # Ensure the file is closed so Camelot can access it

            tables = camelot.read_pdf(temp_file.name, pages="all")
            for table in tables:
                content.append({"type": "table", "data": table.df.values.tolist()})
        except ImportError:
            pass
        except Exception as e:
            print(f"Error processing tables: {e}")
        finally:
            os.remove(temp_file.name)  # Clean up the temporary file
            logger.info("Extraction process is complete")
            logger.info(f"content from extracted pdf: {content[:10]}")
        return content
    
    async def process_excel(self) -> List[Dict]:
        content = []
        xls = pd.ExcelFile(self.document_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if not df.empty:
                content.append({
                    "type": "table",
                    "data": {
                        "sheet_name": sheet_name,
                        "rows": df.fillna("").values.tolist()
                    }
                })
        return content
    
    async def _process_txt(self) -> List[Dict]:
        with open(self.document_path, "r", encoding="utf-8") as f:
            text = f.read().strip()
        return [{"type": "text", "data": text}]
    
    async def _process_pptx(self) -> List[Dict]:
        content = []
        image_count = 0
        prs = Presentation(self.document_path)

        for slide_num, slide in enumerate(prs.slides):
            # Extract text from slide shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append({"type": "text", "data": shape.text.strip()})

                # Extract images
                if shape.shape_type == 13:  # 13 = picture type
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_path = os.path.join("uploads_images", f"{self.document_id}_{self.user_id}_pptx_image_{slide_num+1}_{image_count+1}.{image_ext}")
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    content.append({"type": "image", "data": image_path, "content":ProjectScopingAgent.summarize_image(image_path)})
                    image_count += 1

                # Extract tables
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    content.append({"type": "table", "data": table_data})

        return content

    
    async def extract_csv(self):
        try:
            data = pd.read_csv(filepath_or_buffer=self.document_path)
        except Exception as e:
            return e
        
        return {"data":data}
    
    def parse_document(self):
        if not os.path.exists(self.document_path):
            return "File doesn't exist. Provide valid file path"
        
        file_extension = os.path.splitext(self.document_path)[-1].lower()

        if file_extension == ".docx":
            return self.extract_docx()
        elif file_extension == ".pdf":
            # return self.extract_pdf()
            return self.process_pdf_with_structure()
        elif file_extension == ".txt":
            return self.extract_txt()
        elif file_extension == ".csv":
            return self.extract_csv()
        #removing .xlsx cause of parsing issue for now # too lazy to write parsing in service.py
        # elif file_extension == ".xlsx":
        #     return self.process_excel()
        elif file_extension == ".pptx":
            return self._process_pptx()
        else:
            return "unsupported file type. Please provide .docx .pdf .txt .pptx file"

@asynccontextmanager
async def open_pdf(document_path: str = None, stream: bytes = None, filetype: str = None):
    if document_path:
        doc = fitz.open(document_path)
    elif stream and filetype:
        doc = fitz.open(stream=stream, filetype=filetype)
    else:
        raise ValueError("Either document_path or stream and filetype must be provided")
    try:
        yield doc
    finally:
        doc.close()
        print("Closed the PDF document to conserve resources")

#all authentication is done, want to improve the time time for extracting pdf data and get details of images and table there are 2 pdf functions need to look into it
